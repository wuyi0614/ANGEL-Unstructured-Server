from pptx import Presentation
from unstructured.partition.pptx import partition_pptx


def unstructure_ppt(file_name):
    prs = Presentation(file_name)
    elements = partition_pptx(
        filename=file_name,
        multipage_sections=True,
        infer_table_structure=True,
        include_page_breaks=True,
    )

    text_list = []
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text_frame.text + " "

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text

        slide_elements = elements[i] if i < len(elements) else []
        element_text = []
        if isinstance(slide_elements, list):
            for element in slide_elements:
                if hasattr(element, "text"):
                    element_text.append(element.text)
                else:
                    element_text.append(str(element))
        else:
            element_text.append(str(slide_elements))

        combined_text = f"{slide_text.strip()}\n\n{notes_text.strip()}\n\n{' '.join(element_text)}"

        page_number = i + 1
        text_list.append((combined_text, page_number))

    return text_list
